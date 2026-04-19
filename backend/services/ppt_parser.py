"""
PPT parsing service.
Converts .ppt/.pptx → PDF via LibreOffice headless,
then uses PyMuPDF to extract text and render PNG images per page.
"""

import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

import re

import fitz  # PyMuPDF

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

# Bullet-start patterns: lines beginning with these are standalone bullets
_BULLET_START = re.compile(
    r"^(\s*[•\-\*\u2022\u2023\u25E6\u2043]\s+"    # bullet chars: •, -, *, etc.
    r"|^\s*\d+[\.\)]\s+"                             # numbered: "1. " or "1) "
    r"|^\s*[A-Z][a-z].*[:\.]$"                       # Title-like: starts uppercase, ends : or .
    r")"
)


def _clean_ppt_text(raw_text: str) -> str:
    """
    Post-process PyMuPDF extracted text:
    1. Remove page-number lines (pure digits)
    2. Remove very short standalone lines (<=2 chars)
    3. Remove footer-like lines (1-2 words, no punctuation, e.g. "Xiao Lei")
    4. Merge continuation lines (wrapped bullets) back together
    """
    lines = raw_text.splitlines()
    cleaned: list[str] = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        # Skip pure-digit lines (page numbers like "63")
        if stripped.replace(".", "").isdigit():
            continue
        # Skip very short lines (<=2 chars, likely artifacts)
        if len(stripped) <= 2:
            continue
        # Skip footer-like lines: 1-2 words, no sentence-ending punctuation
        words = stripped.split()
        if (len(words) <= 2
                and not any(c in stripped for c in ".,:;!?()[]{}–—")):
            continue
        cleaned.append(stripped)

    # Merge continuation lines: if a line doesn't look like a new bullet
    # and the previous line doesn't end with sentence-ending punctuation,
    # join it to the previous line.
    merged: list[str] = []
    for line in cleaned:
        is_new_bullet = bool(re.match(
            r"^[•\-\*\u2022\u2023\u25E6\u2043]\s+|^\d+[\.\)]\s+",
            line
        ))
        # Also treat lines starting with an uppercase word followed by more text as new items
        is_title_like = bool(re.match(r"^[A-Z][A-Za-z]+(\s+[A-Z]|\s*:|\s*$)", line))

        if (merged
                and not is_new_bullet
                and not is_title_like
                and not merged[-1].rstrip().endswith((".", ":", ";", "!", "?"))):
            # Continuation of previous line
            merged[-1] = merged[-1].rstrip() + " " + line
        else:
            merged.append(line)

    return "\n".join(merged)


def _libreoffice_path() -> str:
    """Find LibreOffice executable (soffice) on Windows or Linux."""
    candidates = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/lib/libreoffice/program/soffice",
    ]
    for c in candidates:
        if shutil.which(c) or Path(c).exists():
            return c
    raise RuntimeError(
        "LibreOffice (soffice) not found. "
        "Run install_deps.bat and restart your terminal."
    )


def pptx_to_pdf(input_path: str, output_dir: str) -> str:
    """
    Convert a .ppt/.pptx file to PDF using LibreOffice headless.
    Returns the path to the generated PDF file.
    """
    soffice = _libreoffice_path()
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            input_path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed:\n{result.stderr}"
        )

    # LibreOffice names the output file after the input stem
    stem = Path(input_path).stem
    pdf_path = str(Path(output_dir) / f"{stem}.pdf")
    if not Path(pdf_path).exists():
        raise RuntimeError(
            f"Expected PDF not found at {pdf_path}. "
            f"LibreOffice stdout: {result.stdout}"
        )
    return pdf_path


def extract_domain_terms(pages: list[dict], max_terms: int = 50) -> str:
    """
    Extract domain-specific terms from parsed PPT pages for Whisper prompt injection.

    Strategy: collect all words that are likely domain terms —
      - CamelCase words (e.g. BackPropagation, ResNet)
      - ALL_CAPS abbreviations (e.g. CNN, LSTM, API)
      - Words appearing in title-like positions (first word of a bullet)
      - Deduplicated, limited to max_terms

    Returns a comma-separated string suitable for Whisper's `prompt` parameter.
    """
    import re as _re

    camel_or_upper = _re.compile(r"\b([A-Z][a-z]+[A-Z]\w*|[A-Z]{2,}\w*)\b")
    seen: dict[str, int] = {}  # term -> frequency

    for page in pages:
        text = page.get("ppt_text", "") or ""
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            # CamelCase / ALLCAPS terms anywhere on the line
            for match in camel_or_upper.finditer(line):
                term = match.group(1)
                seen[term] = seen.get(term, 0) + 1
            # First "word" of each bullet (often a key concept)
            first_word_match = _re.match(r"^[•\-\*\d\.\)]*\s*([A-Za-z][A-Za-z\-]{2,})", line)
            if first_word_match:
                term = first_word_match.group(1)
                if term[0].isupper():
                    seen[term] = seen.get(term, 0) + 1

    # Sort by frequency descending, take top max_terms
    top_terms = sorted(seen, key=lambda t: seen[t], reverse=True)[:max_terms]
    return ", ".join(top_terms)


def _extract_text_pptx(pptx_path: str) -> list[str]:
    """
    Extract text per slide from a .pptx file using python-pptx.

    Each slide's text is structured as one line per paragraph (preserving the
    original bullet / title structure). Empty slides return an empty string.

    Returns a list of strings, one per slide (0-indexed → slide 1 is index 0).
    """
    if not _HAS_PPTX:
        raise RuntimeError(
            "python-pptx is required for .pptx text extraction. "
            "Install it: pip install python-pptx"
        )

    prs = _Presentation(pptx_path)
    slide_texts: list[str] = []

    for slide in prs.slides:
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Skip pure page-number lines
                if text.replace(".", "").isdigit():
                    continue
                lines.append(text)
        slide_texts.append("\n".join(lines))

    return slide_texts


def _extract_bullets_pptx(pptx_path: str) -> list[list[dict]]:
    """
    Extract structured bullets per slide with level information from a .pptx file.

    Returns a list (one per slide) of bullet lists, each bullet being:
      {
        "text": str,    # paragraph text
        "level": int,   # 0 = title/heading, 1 = first-level bullet, 2 = sub-bullet, etc.
      }

    Heuristic: shapes are sorted so title placeholders come first (level 0),
    then body text frames use paragraph.level directly.
    """
    if not _HAS_PPTX:
        return []

    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    prs = _Presentation(pptx_path)
    slides_bullets: list[list[dict]] = []

    for slide in prs.slides:
        bullets: list[dict] = []

        # Separate title placeholders from body shapes so titles always come first
        title_shapes = []
        body_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.idx == 0:
                title_shapes.append(shape)
            else:
                body_shapes.append(shape)

        for shape in title_shapes:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text or text.replace(".", "").isdigit():
                    continue
                bullets.append({"text": text, "level": 0})

        for shape in body_shapes:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text or text.replace(".", "").isdigit():
                    continue
                # python-pptx paragraph.level: 0 = top-level body, 1 = first indent, etc.
                # We add 1 so that body level 0 becomes level 1 (below the title at level 0)
                level = (para.level or 0) + 1
                bullets.append({"text": text, "level": level})

        slides_bullets.append(bullets)

    return slides_bullets


def parse_ppt(
    ppt_path: str,
    slides_output_dir: str,
    pdf_name: str = "slides.pdf",
) -> list[dict]:
    """
    Full PPT parsing pipeline:
      1. Convert PPT/PPTX to PDF via LibreOffice (for rendering)
      2. Copy PDF to slides_output_dir
      3. Extract text:
         - .pptx: use python-pptx (accurate paragraph structure)
         - .pdf:  use PyMuPDF (fallback)
         - .ppt:  LibreOffice converts to .pptx first, then python-pptx
      4. For each PDF page: render PNG thumbnail

    Returns a list of dicts:
      {
        "page_num": int,           # 1-based
        "ppt_text": str,           # extracted text (for semantic alignment)
        "pdf_url": str,            # relative URL to the PDF
        "pdf_page_num": int,       # 1-based page index within the PDF
        "thumbnail_url": str,      # relative URL to the PNG thumbnail
      }
    """
    os.makedirs(slides_output_dir, exist_ok=True)

    suffix = Path(ppt_path).suffix.lower()

    with tempfile.TemporaryDirectory() as tmp_dir:
        # --- Determine text extraction source ---
        pptx_texts: list[str] | None = None
        pptx_bullets: list[list[dict]] | None = None

        pptx_bullets: list[list[dict]] | None = None

        if suffix == ".pptx" and _HAS_PPTX:
            # Direct python-pptx extraction from the original file
            pptx_texts = _extract_text_pptx(ppt_path)
            pptx_bullets = _extract_bullets_pptx(ppt_path)

        elif suffix == ".ppt" and _HAS_PPTX:
            # Convert .ppt → .pptx via LibreOffice, then extract with python-pptx
            soffice = _libreoffice_path()
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pptx",
                 "--outdir", tmp_dir, ppt_path],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0:
                pptx_file = str(Path(tmp_dir) / (Path(ppt_path).stem + ".pptx"))
                if Path(pptx_file).exists():
                    pptx_texts = _extract_text_pptx(pptx_file)
                    pptx_bullets = _extract_bullets_pptx(pptx_file)

        # --- Convert to PDF for rendering ---
        if suffix == ".pdf":
            src_pdf = ppt_path
        else:
            src_pdf = pptx_to_pdf(ppt_path, tmp_dir)

        # Copy PDF into slides_output_dir for serving
        dest_pdf = Path(slides_output_dir) / pdf_name
        shutil.copy2(src_pdf, dest_pdf)
        pdf_url = f"/slides/{pdf_name}"

        # --- Per-page: text + PNG thumbnail ---
        doc = fitz.open(src_pdf)
        pages = []

        for i, pdf_page in enumerate(doc):
            page_num = i + 1

            # Text: prefer python-pptx extraction, fallback to PyMuPDF
            if pptx_texts is not None and i < len(pptx_texts):
                ppt_text = pptx_texts[i]
            else:
                ppt_text = _clean_ppt_text(pdf_page.get_text("text"))

            # No pre-rendering: frontend falls back to on-demand /api/sessions/{id}/slide/{n}.png
            thumbnail_url = None

            pages.append(
                {
                    "page_num": page_num,
                    "ppt_text": ppt_text,
                    "ppt_bullets": (
                        pptx_bullets[i]
                        if pptx_bullets is not None and i < len(pptx_bullets)
                        else None
                    ),
                    "pdf_url": pdf_url,
                    "pdf_page_num": page_num,
                    "thumbnail_url": thumbnail_url,
                }
            )

        doc.close()

    return pages
